#!/usr/bin/env python3
"""
veeva_extract.py

Standalone script to extract text from files (PDF, DOCX, PPTX, XLSX, images).
- Native extraction first.
- Selective OCR on pages with low native text coverage.
- Image OCR merging into page text, plus separate per-page OCR image text buckets.
- Chunking helper for downstream LLM queries.

Designed for Linux EMR server. Tune concurrency/thresholds at top.
"""

import os
import re
import io
import sys
import math
import json
import logging
import tempfile
import multiprocessing
from functools import partial
from concurrent.futures import ProcessPoolExecutor, as_completed

# file handlers
import fitz  # PyMuPDF
import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd

# OCR & image processing
from PIL import Image, ImageFilter, ImageOps
import pytesseract
import cv2
import numpy as np

# -----------------------
# Configuration (tune)
# -----------------------
OCR_TEXT_COVERAGE_THRESHOLD = 0.30  # if native text length / expected_page_capacity < this -> OCR page
MIN_NATIVE_CHARS_TO_SKIP_OCR = 120  # quick heuristic: if chars >= this, skip page OCR
CHUNK_MAX_CHARS = 6000
CHUNK_OVERLAP = 300
MAX_OCR_WORKERS = min(multiprocessing.cpu_count(), 4)  # safe default for EMR; increase if you know resources
TESSERACT_LANG = "eng"  # change to "eng+fra" etc if needed
TESSERACT_PSM = 3  # 3=Fully automatic page segmentation
TESSERACT_OEM = 1  # LSTM only
TESSERACT_MIN_CONFIDENCE = 45.0  # below this mean confidence -> consider fallback or flag
USE_OCRMYPDF_FALLBACK = False  # set True to run ocrmypdf for scanned PDFs as fallback
LOG_LEVEL = logging.INFO

# Expected page text capacity heuristic (used to compute coverage)
# For typical A4 page with dense text, assume 3000 characters. Adjust if your docs are mostly slides (lower).
EXPECTED_PAGE_TEXT_CAPACITY = 3000

# -----------------------
# Logging
# -----------------------
logging.basicConfig(level=LOG_LEVEL, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger("veeva_extract")

# -----------------------
# Utilities
# -----------------------
def normalize_text(s: str) -> str:
    if not s:
        return ""
    # normalize whitespace and ligatures
    s = s.replace("\ufb01", "fi").replace("\ufb02", "fl")
    s = re.sub(r"\s+", " ", s).strip()
    return s

def sentence_split(text):
    # naive sentence split; fine for chunking boundaries
    return re.split(r'(?<=[.!?])\s+', text)

def chunk_text_by_chars(text, max_chars=CHUNK_MAX_CHARS, overlap=CHUNK_OVERLAP):
    text = text or ""
    if len(text) <= max_chars:
        return [text]
    chunks = []
    start = 0
    L = len(text)
    while start < L:
        end = start + max_chars
        if end >= L:
            chunks.append(text[start:L])
            break
        # try to move end back to last sentence boundary within 300 chars
        snippet = text[start:end]
        m = re.search(r'(?<=\S)(?=[\.\?\!]\s)', snippet[::-1])  # not reliable reversed; fallback below
        # simpler: find last period in last 200 chars
        back = snippet.rfind(".", max(0, len(snippet)-200))
        if back != -1:
            end = start + back + 1
        chunk = text[start:end]
        chunks.append(chunk)
        start = max(0, end - overlap)
    return chunks

# -----------------------
# Image preprocessing for OCR
# -----------------------
def preprocess_pil_for_ocr(pil_img: Image.Image, target_dpi=300):
    """
    Preprocess PIL image for better OCR accuracy:
    - convert to grayscale
    - upscale to target DPI if small
    - adaptive threshold
    - denoise
    - deskew
    """
    # convert to RGB then grayscale
    img = pil_img.convert("RGB")
    w, h = img.size

    # upscale if small so OCR sees enough pixels
    # heuristics: ensure smallest dimension >= 1000 px
    min_dim = min(w, h)
    if min_dim < 1000:
        scale = max(1, 1000 // min_dim)
        new_w, new_h = (int(w * scale), int(h * scale))
        img = img.resize((new_w, new_h), resample=Image.BICUBIC)

    # convert to grayscale numpy array
    arr = np.array(img.convert("L"))

    # denoise
    arr = cv2.medianBlur(arr, 3)

    # adaptive threshold
    arr = cv2.adaptiveThreshold(arr, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                cv2.THRESH_BINARY, 31, 15)

    # deskew using moments
    coords = np.column_stack(np.where(arr > 0))
    angle = 0.0
    try:
        rect = cv2.minAreaRect(coords)
        angle = rect[-1]
        if angle < -45:
            angle = -(90 + angle)
        else:
            angle = -angle
        # rotate image to deskew
        (h2, w2) = arr.shape[:2]
        center = (w2 // 2, h2 // 2)
        M = cv2.getRotationMatrix2D(center, angle, 1.0)
        arr = cv2.warpAffine(arr, M, (w2, h2), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)
    except Exception:
        # if deskew fails, continue with original
        pass

    proc_img = Image.fromarray(arr)
    return proc_img

def pytesseract_ocr_with_confidence(pil_img: Image.Image, lang=TESSERACT_LANG, psm=TESSERACT_PSM, oem=TESSERACT_OEM):
    """
    Run tesseract, return (text, mean_confidence, word_confidences)
    """
    config = f"--psm {psm} --oem {oem}"
    try:
        # get TSV output for confidences
        tsv = pytesseract.image_to_data(pil_img, lang=lang, config=config, output_type=pytesseract.Output.DICT)
    except Exception as e:
        logger.exception("Tesseract OCR failure: %s", e)
        return "", 0.0, []

    words = []
    confs = []
    text_parts = []
    n = len(tsv.get("text", []))
    for i in range(n):
        w = tsv["text"][i]
        conf = float(tsv["conf"][i]) if tsv["conf"][i].strip() != "-1" and tsv["conf"][i] != "" else -1.0
        if w.strip():
            words.append(w)
            confs.append(conf)
            text_parts.append(w)
    full_text = " ".join(text_parts)
    mean_conf = float(np.mean([c for c in confs if c >= 0]) ) if confs else 0.0
    return full_text, mean_conf, confs

# -----------------------
# File extractors
# -----------------------
def extract_text_from_docx(path):
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text]
    return " ".join(paragraphs)

def extract_text_from_pptx(path):
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        slides_text.append("\n".join(texts))
    # join slides as pages
    return slides_text

def extract_text_from_xlsx(path):
    xls = pd.ExcelFile(path)
    sheets_text = {}
    for sheet in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet, dtype=str)
        # flatten rows to text
        rows = []
        for r in df.fillna("").astype(str).values:
            rows.append("\t".join(r))
        sheets_text[sheet] = "\n".join(rows)
    return sheets_text

# -----------------------
# PDF helpers
# -----------------------
def extract_pdf_native_pages(pdf_path):
    """
    Returns list of native text strings per page (empty string if none).
    """
    doc = fitz.open(pdf_path)
    pages = []
    for page in doc:
        text = page.get_text("text") or ""
        pages.append(text)
    return pages

def render_pdf_page_to_pil(page, zoom=2.0):
    """
    Render fitz page to PIL Image.
    zoom=2.0 ~ approximates ~144 DPI. Increase for higher DPI.
    """
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img_bytes = pix.tobytes("png")
    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    return img

def extract_images_from_pypdf_page(page):
    """
    Extract embedded images from a PyMuPDF page (returns list of PIL images).
    """
    images = []
    for img_index, img in enumerate(page.get_images(full=True)):
        xref = img[0]
        base_image = page.get_image_stream(xref)
        try:
            pil = Image.open(io.BytesIO(base_image)).convert("RGB")
            images.append(pil)
        except Exception:
            continue
    return images

# -----------------------
# Page processing orchestration
# -----------------------
def should_ocr_page(native_text: str):
    """
    Decide whether to OCR page: if native_text coverage low.
    """
    native_len = len(native_text.strip()) if native_text else 0
    coverage = (native_len / EXPECTED_PAGE_TEXT_CAPACITY) if EXPECTED_PAGE_TEXT_CAPACITY > 0 else 0.0
    if native_len >= MIN_NATIVE_CHARS_TO_SKIP_OCR:
        return False
    return coverage < OCR_TEXT_COVERAGE_THRESHOLD

def process_pdf_page_for_ocr(page_tuple):
    """
    Worker-friendly function expecting (pdf_path, page_number).
    Returns dict: {page: i, native_text:..., ocr_images_text: [...], ocr_text: ..., ocr_mean_conf: float}
    """
    pdf_path, page_index = page_tuple
    doc = fitz.open(pdf_path)
    page = doc[page_index]
    native_text = page.get_text("text") or ""
    # If native text looks sufficient and threshold not met, we skip heavy OCR
    do_ocr = should_ocr_page(native_text)
    ocr_images_text = []
    ocr_text = ""
    ocr_mean_conf = 0.0
    if do_ocr:
        # render page to image at higher zoom for OCR
        pil_img = render_pdf_page_to_pil(page, zoom=3.0)  # higher DPI for accuracy
        proc = preprocess_pil_for_ocr(pil_img)
        txt, mean_conf, confs = pytesseract_ocr_with_confidence(proc)
        ocr_text = normalize_text(txt)
        ocr_mean_conf = mean_conf
        # Also extract embedded images (if any) and OCR them separately (rare)
        try:
            imgs = extract_images_from_pypdf_page(page)
            for img in imgs:
                p = preprocess_pil_for_ocr(img)
                t, c, cs = pytesseract_ocr_with_confidence(p)
                if t.strip():
                    ocr_images_text.append(normalize_text(t))
        except Exception:
            pass
    else:
        # If not doing full-page OCR, we may still OCR embedded images if they exist AND native text is missing from them.
        try:
            imgs = extract_images_from_pypdf_page(page)
            for img in imgs:
                p = preprocess_pil_for_ocr(img)
                t, c, cs = pytesseract_ocr_with_confidence(p)
                if t.strip():
                    ocr_images_text.append(normalize_text(t))
        except Exception:
            pass

    return {
        "page": page_index + 1,
        "native_text": normalize_text(native_text),
        "ocr_text": ocr_text,
        "ocr_images_text": ocr_images_text,
        "ocr_mean_conf": ocr_mean_conf,
        "did_ocr": do_ocr
    }

# -----------------------
# Top-level processors
# -----------------------
def process_pdf_file(path, workers=MAX_OCR_WORKERS):
    logger.info("Processing PDF: %s", path)
    doc = fitz.open(path)
    n_pages = doc.page_count
    # First pass: native extraction per page
    native_pages = [doc[i].get_text("text") or "" for i in range(n_pages)]

    # decide pages that need full OCR
    tasks = []
    for i, native in enumerate(native_pages):
        if should_ocr_page(native):
            tasks.append((path, i))
        else:
            # we will still attempt to extract embedded images text later in worker run
            tasks.append((path, i))  # we pass all pages to worker which will be lightweight for non-ocr pages

    results = [None] * n_pages
    # Use process pool for OCR (since tesseract uses C libs and benefits from processes)
    with ProcessPoolExecutor(max_workers=workers) as exc:
        futures = {exc.submit(process_pdf_page_for_ocr, t): t for t in tasks}
        for fut in as_completed(futures):
            try:
                res = fut.result()
                idx = res["page"] - 1
                # merge: native_text + ocr_text + ocr_images_text merged on page
                native = res["native_text"] or ""
                merged = native
                # prefer native text, append ocr_text if native is empty or short
                if res["ocr_text"]:
                    if not native.strip():
                        merged = res["ocr_text"]
                    else:
                        # append OCR image text or OCR fallback text if distinct
                        if res["ocr_text"].strip() and res["ocr_text"].strip() not in native:
                            merged = merged + "\n " + res["ocr_text"]
                # append texts recognized from images
                if res["ocr_images_text"]:
                    merged = merged + "\n " + "\n ".join(res["ocr_images_text"])
                results[idx] = {
                    "page": res["page"],
                    "native_text": native,
                    "ocr_text": res["ocr_text"],
                    "ocr_images_text": res["ocr_images_text"],
                    "merged_text": normalize_text(merged),
                    "ocr_mean_conf": res["ocr_mean_conf"],
                    "did_ocr": res["did_ocr"]
                }
            except Exception as e:
                logger.exception("Error processing page task: %s", e)

    # fallback: if many pages have very low OCR confidence and USE_OCRMYPDF_FALLBACK, try full document ocrmypdf
    if USE_OCRMYPDF_FALLBACK:
        low_conf_pages = [r for r in results if r and r.get("ocr_mean_conf", 100.0) < TESSERACT_MIN_CONFIDENCE]
        if len(low_conf_pages) > 0:
            logger.warning("Low confidence detected on %d pages; consider running ocrmypdf fallback.", len(low_conf_pages))
            # One could run ocrmypdf here on the full file and rerun extraction (left as optional)

    # Build final outputs
    pages_struct = results
    full_text = "\n\n".join([p["merged_text"] for p in pages_struct if p])
    return {"full_text": normalize_text(full_text), "pages": pages_struct}

def process_docx_file(path):
    txt = extract_text_from_docx(path)
    merged = normalize_text(txt)
    return {"full_text": merged, "pages": [{"page": 1, "native_text": merged, "ocr_text": "", "ocr_images_text": [], "merged_text": merged}]}

def process_pptx_file(path):
    slides = extract_text_from_pptx(path)
    pages = []
    for i, s in enumerate(slides):
        native = normalize_text(s)
        # slides may have images; to keep accuracy, we could extract images and OCR them (not implemented by default)
        pages.append({"page": i+1, "native_text": native, "ocr_text": "", "ocr_images_text": [], "merged_text": native})
    full_text = "\n\n".join([p["merged_text"] for p in pages])
    return {"full_text": normalize_text(full_text), "pages": pages}

def process_xlsx_file(path):
    sheets = extract_text_from_xlsx(path)
    pages = []
    for i, (sheetname, text) in enumerate(sheets.items()):
        native = normalize_text(text)
        pages.append({"page": i+1, "sheet": sheetname, "native_text": native, "ocr_text": "", "ocr_images_text": [], "merged_text": native})
    full_text = "\n\n".join([p["merged_text"] for p in pages])
    return {"full_text": normalize_text(full_text), "pages": pages}

def process_image_file(path):
    img = Image.open(path)
    proc = preprocess_pil_for_ocr(img)
    txt, mean_conf, confs = pytesseract_ocr_with_confidence(proc)
    merged = normalize_text(txt)
    return {"full_text": merged, "pages": [{"page": 1, "native_text": "", "ocr_text": merged, "ocr_images_text": [], "merged_text": merged, "ocr_mean_conf": mean_conf}]}

# master dispatcher
def process_file(path, workers=MAX_OCR_WORKERS):
    path = os.path.abspath(path)
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return process_pdf_file(path, workers=workers)
    elif ext == ".docx":
        return process_docx_file(path)
    elif ext in (".pptx",):
        return process_pptx_file(path)
    elif ext in (".xls", ".xlsx"):
        return process_xlsx_file(path)
    elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
        return process_image_file(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

# -----------------------
# CLI and sample usage
# -----------------------
def main():
    import argparse
    parser = argparse.ArgumentParser(description="Extract text from Veeva files (PDF/DOCX/PPTX/XLSX/images).")
    parser.add_argument("input", help="Path to file or directory")
    parser.add_argument("--workers", type=int, default=MAX_OCR_WORKERS, help="Number of OCR worker processes")
    parser.add_argument("--outdir", default="veeva_out", help="Output directory for extracted JSON/text")
    parser.add_argument("--chunk", action="store_true", help="Also output LLM-ready chunks")
    args = parser.parse_args()

    inp = args.input
    outdir = args.outdir
    os.makedirs(outdir, exist_ok=True)

    targets = []
    if os.path.isdir(inp):
        for fn in os.listdir(inp):
            fp = os.path.join(inp, fn)
            if os.path.isfile(fp):
                targets.append(fp)
    else:
        targets = [inp]

    summary = {}
    for f in targets:
        try:
            logger.info("Processing %s", f)
            res = process_file(f, workers=args.workers)
            base = os.path.basename(f)
            name = os.path.splitext(base)[0]
            out_json = os.path.join(outdir, f"{name}.json")
            with open(out_json, "w", encoding="utf-8") as fh:
                json.dump(res, fh, ensure_ascii=False, indent=2)
            # also write full text
            txt_path = os.path.join(outdir, f"{name}.txt")
            with open(txt_path, "w", encoding="utf-8") as fh:
                fh.write(res["full_text"] or "")
            # optional chunk output
            if args.chunk:
                chunks = chunk_text_by_chars(res["full_text"], max_chars=CHUNK_MAX_CHARS, overlap=CHUNK_OVERLAP)
                chunks_path = os.path.join(outdir, f"{name}.chunks.json")
                with open(chunks_path, "w", encoding="utf-8") as fh:
                    json.dump({"chunks": chunks}, fh, ensure_ascii=False, indent=2)
            summary[f] = {"out_json": out_json, "txt": txt_path}
            logger.info("Wrote outputs for %s -> %s", f, out_json)
        except Exception as e:
            logger.exception("Failed to process %s: %s", f, e)
    # write summary
    with open(os.path.join(outdir, "summary.json"), "w", encoding="utf-8") as fh:
        json.dump(summary, fh, ensure_ascii=False, indent=2)
    logger.info("Done. Outputs in %s", outdir)

if __name__ == "__main__":
    main()
